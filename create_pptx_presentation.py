# -*- coding: utf-8 -*-
"""
YES24 IT/모바일 베스트셀러 EDA 분석 결과를 바탕으로 
노르딕 미니멀리즘 스타일의 30페이지 분량 PPTX 프레젠테이션을 생성하는 파이썬 스크립트입니다.
제목은 Arial, 본문은 나눔고딕으로 설정하고 슬라이드 노트를 포함하며 인포그래픽 도형을 배치합니다.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_presentation(output_filepath):
    print("PPTX 프레젠테이션 생성 시작...")
    prs = Presentation()
    
    # 16:9 비율 설정 (가로 13.333인치, 세로 7.5인치)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # 빈 슬라이드 레이아웃
    
    # --- 테마 색상 정의 (노르딕 미니멀리즘 - 벽돌색 포인트) ---
    COLOR_BG = RGBColor(0xFA, 0xF9, 0xF6)       # 미색 (Warm Off-white)
    COLOR_TEXT = RGBColor(0x2F, 0x3E, 0x46)     # 어두운 차콜 (Charcoal)
    COLOR_TITLE = RGBColor(0xB8, 0x50, 0x42)    # 포인트 컬러: 벽돌색 (Brick Red / Terracotta)
    COLOR_ACCENT = RGBColor(0xB8, 0x50, 0x42)   # 벽돌색 (Brick Red)
    COLOR_CARD = RGBColor(0xEE, 0xEB, 0xE4)     # 노르딕 베이지 (Nordic Beige)
    COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    COLOR_LEAD_BG = RGBColor(0xB8, 0x50, 0x42)  # 벽돌색 간지 배경
    COLOR_LEAD_TEXT = RGBColor(0xFA, 0xF9, 0xF6)
    
    # --- 헬퍼 함수 정의 ---
    
    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_flat_icon(slide, x, y, size, icon_type):
        # 원형 배경
        bg_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        bg_circle.fill.solid()
        bg_circle.fill.fore_color.rgb = COLOR_TITLE
        bg_circle.line.fill.background()
        
        # 내부 기호 배치 (흰색으로 대비)
        inner_size = size * 0.5
        offset = (size - inner_size) / 2
        ix, iy = x + offset, y + offset
        
        if icon_type == "arrow":
            shape_type = MSO_SHAPE.RIGHT_ARROW
        elif icon_type == "star":
            shape_type = MSO_SHAPE.STAR_5_POINT
        elif icon_type == "smile":
            shape_type = MSO_SHAPE.SMILEY_FACE
        elif icon_type == "cloud":
            shape_type = MSO_SHAPE.CLOUD
        elif icon_type == "diamond":
            shape_type = MSO_SHAPE.DIAMOND
        elif icon_type == "heart":
            shape_type = MSO_SHAPE.HEART
        else:
            shape_type = MSO_SHAPE.RECTANGLE
            
        icon = slide.shapes.add_shape(shape_type, Inches(ix), Inches(iy), Inches(inner_size), Inches(inner_size))
        icon.fill.solid()
        icon.fill.fore_color.rgb = COLOR_BG
        icon.line.fill.background()
        return bg_circle

    def add_title(slide, text):
        # 제목 텍스트박스
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = 0
        tf.margin_left = 0
        
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = 'Arial'
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = COLOR_TITLE
        
        # 제목 밑의 얇은 장식 선 (노르딕 스타일)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.fill.background() # 선 없음

    def add_card(slide, x, y, w, h, fill_color, border_color=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    def add_text_box(slide, x, y, w, h, text_runs, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.05)
        tf.margin_left = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.margin_right = Inches(0.05)
        
        for idx, run_data in enumerate(text_runs):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.alignment = align
            
            # 불릿 여부
            is_bullet = run_data.get('bullet', False)
            if is_bullet:
                p.level = 0
                
            text = run_data.get('text', '')
            font_name = run_data.get('font', '나눔고딕')
            size = run_data.get('size', 14)
            color = run_data.get('color', COLOR_TEXT)
            bold = run_data.get('bold', False)
            italic = run_data.get('italic', False)
            
            run = p.add_run()
            run.text = text
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color
            
        return txBox

    def add_image_slide_content(slide, img_filename, md_dir, title_text, desc_text_runs):
        add_title(slide, title_text)
        
        # 이미지 경로 빌드
        img_path = os.path.normpath(os.path.join(md_dir, img_filename))
        
        # 좌측 이미지 배치
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.5))
        else:
            # 이미지 부재 시 대체 카드
            card = add_card(slide, 0.8, 1.8, 6.0, 4.5, COLOR_CARD)
            tx = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(5.6), Inches(1.0))
            p = tx.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = f"[이미지를 찾을 수 없음: {img_filename}]"
            r.font.name = "나눔고딕"
            r.font.size = Pt(16)
            r.font.color.rgb = RGBColor(180, 50, 50)
            
        # 우측 해설 카드 배치
        add_card(slide, 7.2, 1.8, 5.3, 4.5, COLOR_CARD)
        add_text_box(slide, 7.4, 2.0, 4.9, 4.1, desc_text_runs)

    def create_lead_slide(title, subtitle):
        slide = prs.slides.add_slide(blank_layout)
        set_background(slide, COLOR_LEAD_BG)
        
        # 세로 중앙 정렬을 위한 텍스트박스
        txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = 'Arial'
        run.font.size = Pt(44)
        run.font.bold = True
        run.font.color.rgb = COLOR_LEAD_TEXT
        
        if subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(14)
            run2 = p2.add_run()
            run2.text = subtitle
            run2.font.name = '나눔고딕'
            run2.font.size = Pt(22)
            run2.font.color.rgb = RGBColor(0xD9, 0xE1, 0xF2)
            
        return slide

    # ==========================================
    # 발표자 노트 대본 리스트 (30개 장표)
    # ==========================================
    scripts = [
        # 1. 표지
        "안녕하십니까, 오늘 발표를 맡은 데이터 분석 연구원입니다. 이번 세션에서는 YES24 IT/모바일 베스트셀러 도서 데이터를 탐색적으로 분석하여 데이터 내부에 숨겨진 가격 정책, 판매지수의 흥행 역학, 출판사별 효율성, 주요 기술 키워드 및 독자들의 평점 패턴에 대해 상세히 다루겠습니다. 869건의 완전한 베스트셀러 도서 데이터를 정량적으로 정제하여 구축한 만큼, IT 서적 시장의 비즈니스적 통찰을 도출하는 데 훌륭한 나침반이 될 것으로 기대합니다. 차분하면서도 가독성 높은 북유럽풍 미니멀리즘 스타일로 디자인된 본 분석 모델의 전개 과정을 즐겁게 따라와 주시기 바랍니다.",
        
        # 2. 목차
        "본 발표의 전체 목차를 설명해 드리겠습니다. 크게 여덟 가지 영역으로 흐름이 유기적으로 연결됩니다. 첫째, 분석의 배경과 거시적 통계를 파악하고, 둘째, 도서 가격대의 집중도와 그 해석에 대해 알아봅니다. 셋째, 판매 성과를 자극하는 변수 간 상관관계를 다루며, 넷째, 시장의 메이저 출판사들이 어떤 효율로 브랜드를 운영 중인지 비교합니다. 다섯째, 독자 반응의 품질 마지노선이라 할 수 있는 리뷰 평점 패턴을 짚어보고, 여섯째, 짧아진 IT 서적 소비 연도를 실증합니다. 일곱째, TF-IDF 핵심 키워드 30선을 정량화하여 기획 트렌드를 파악하고 마지막으로 비즈니스 액션 플랜을 제시하겠습니다.",
        
        # 3. [간지] 개요 및 기술통계
        "첫 번째 파트인 분석 개요 및 기초 통계 세션을 시작하겠습니다. 우리는 본 세션에서 데이터의 수집 및 정제 방식부터, 베스트셀러 시장 전체의 거시적 요약 지표, 그리고 개별 변수들의 평균, 최솟값, 최댓값, 중앙값 등의 기본 통계 분포를 세밀하게 들여다볼 예정입니다. 데이터의 거시적 형태를 잘 이해해야만 후속으로 분석할 인과관계 통계가 신뢰성을 얻기 때문입니다.",
        
        # 4. 분석 배경 및 데이터 소개
        "본 분석의 구체적인 목적과 데이터셋의 수집 환경을 소개해 드리겠습니다. 현재 IT 서적 소비 시장은 기술 패러다임이 시시각각 변화하고 있습니다. 이에 따라 도서 수집 리스트를 구축했고, YES24 IT/모바일 분야의 종합 베스트셀러 총 869건의 다차원 데이터를 정밀 검수했습니다. 상품 식별 정보부터 정가, 판매가, 판매지수, 리뷰 건수, 리뷰 평점 및 출판일 등 총 15개의 복합 필드를 수집하여 상관계수 및 기술통계 연산의 기초 체력을 마련했습니다. 이번 정량 연구를 통해 향후 출판의 최적 포지셔닝 가격 설정과 마케팅의 리소스 분배 기준을 세우고자 합니다.",
        
        # 5. 데이터 전처리 및 정제
        "데이터의 통계적 타당성을 위해 거친 정밀 정제 단계를 요약해 드리겠습니다. 저자명이 누락된 2건에 대해 통계 왜곡을 피하기 위해 '저자 미상'으로 즉각 대체했고, 부제목의 결측치는 있는 그대로 비워 두어 실무서의 부제목 미적용 비중을 확보했습니다. 특히 '포인트' 열은 콤마와 한글 단어가 섞여 있던 문자를 파싱하여 온전한 정수값으로 정제했고, 10% 단위의 할인율은 엑셀과 파이썬 연산의 호환성을 극대화하기 위해 소수로 변환했습니다. 중복 검출 테스트 결과 0건의 중복을 기록하여 완벽한 분석 데이터셋을 완성했습니다.",
        
        # 6. 거시적 주요 요약 지표 (KPI)
        "베스트셀러 전체 시장의 거시적 실적 지표, 즉 핵심 요약 KPI를 설명드리겠습니다. 총 869권 도서들의 평균 정가는 25,560원이며 실제 구매 가격은 평균 23,395원으로 기본 8.4%의 할인율이 균등하게 작동하고 있습니다. 이는 도서정가제의 영향 아래에서 대다수의 도서들이 10% 고정 할인을 주는 시장의 경직성을 대변합니다. 베스트셀러들의 누적 판매지수 총합은 약 235만 건에 이르며, 독자들의 평균적인 평점은 10점 만점에 7.55점으로 전반적으로 우호적인 독자 만족도를 형성하고 있음을 데이터로 알 수 있습니다.",
        
        # 7. 수치 데이터 기술통계 요약
        "화면의 기술통계 요약 테이블은 엑셀 수식을 통해 자동으로 집계된 결과입니다. 가격, 할인율, 포인트, 판매지수, 리뷰 건수 및 평점의 기초적 기술통계량을 보여줍니다. 정가의 중앙값은 24,400원이며 평균 정가는 25,560원으로 비교적 대칭적이고 안정적인 단가 분포를 따릅니다. 반면, 흥행 실적을 뜻하는 판매지수는 최소 60부터 최대 약 8만 점까지 표준편차가 엄청난 격차를 보여주어 개별 도서의 시장 흥행 실적이 결코 균등하지 않으며 매우 소수에 집중되고 있음을 통계적으로 보여줍니다.",
        
        # 8. 기술통계에서 발견된 비즈니스 롱테일
        "기초 통계에서 포착된 비즈니스 롱테일 현상을 비즈니스적으로 심층 분석해보겠습니다. 판매지수의 산술 평균은 2,701점인 반면, 정중앙에 위치한 도서의 지수는 1,062점에 불과합니다. 평균과 중앙값의 격차가 무려 2.5배에 이른다는 통계적 사실은, 극소수의 베스트셀러 도서가 판매 실적 전체의 대부분을 독점하고 있음을 의미합니다. 리뷰 건수 역시 평균 18.6건 대비 중앙값은 9건으로, 전형적인 파레토 법칙이 IT 기술서 시장에서도 그대로 증명됩니다. 따라서 자원을 분산 투자하기보단 킬러 타이틀에 집중 마케팅하는 기획이 중요합니다.",
        
        # 9. [간지] 도서 가격대 분포 분석
        "두 번째 파트인 도서 가격대 분포 분석 세션을 시작하겠습니다. 본 세션에서는 가격 분포 히스토그램을 실제로 확인하고, IT/모바일 도서의 가격대별 비즈니스 특징과 저항선을 짚어보겠습니다. 이를 통해 출판 마진 확보와 고객의 구매 장벽을 극복할 수 있는 최적의 가격 가이드라인을 제안해 드리고자 합니다.",
        
        # 10. 도서 가격 분포 시각화
        "화면에 표시된 히스토그램은 도서 가격의 분포 범위를 시각적으로 나타낸 것입니다. 보시는 바와 같이 IT/모바일 서적의 단가는 15,000원 이하의 저가 영역보다는, 주로 15,000원부터 30,000원 사이의 높은 정가 포지션에 아주 두껍고 밀집된 히스토그램 피크를 형성하고 있음을 확인할 수 있습니다. 3만 원 이상의 고단가 도서들 역시 넓은 꼬리를 이루며 지속적으로 유통되고 있는데, 이는 IT 서적 특유의 전문성과 개발 실습 예제의 방대함이 도서 제조 단가 및 공급 가격에 직접적으로 투영된 실증 지표입니다.",
        
        # 11. 가격대 구간별 분포의 비즈니스적 해석
        "가격대 구간별 비즈니스 시사점을 세부적으로 풀어드리겠습니다. 첫째, 15,000원 이하의 저단가 구간은 페이지 수가 적은 기초 모바일 기기 매뉴얼이나 기출문제 위주의 수험서, 혹은 전자책 비중이 높습니다. 둘째, 15,000원~30,000원의 표준 실무 서적 구간은 파이썬 입문서, 웹 프레임워크 책 등 캐시카우 역할을 톡톡히 하는 볼륨 구간군군입니다. 마지막으로 35,000원 이상의 고가 전문 서적 구간은 AI 대학 전공서나 클라우드 설계 가이드처럼 기술 난이도가 높고 두께가 두꺼운 하드웨어/학술 전문서가 포진하여 출판사 매출 단가를 견인합니다.",
        
        # 12. 최적 가격 포지셔닝 및 기획 전략
        "따라서 당사가 신간을 기획할 때 어떠한 가격 정책을 세워야 할지에 대한 명확한 결론을 내리겠습니다. IT 독자들은 책 가격 2,000~3,000원 차이보다 콘텐츠의 깊이와 예제 퀄리티를 최우선으로 고려하는 고관여 타겟군입니다. 그러므로 섣부른 저가 경쟁을 위해 페이지 수를 무리하게 축소하여 품질을 낮추는 것은 자폭에 가깝습니다. 차라리 페이지 볼륨과 풍부한 삽화, 최신 API를 확실하게 보장하고 **24,000원~28,000원**의 표준 실무 가격대를 책정하여 마진을 확고히 챙기는 기획 전략이 비즈니스적으로 훨씬 똑똑합니다.",
        
        # 13. [간지] 변수 간 상관관계 분석
        "세 번째 파트인 변수 간 상관관계 분석 세션을 시작하겠습니다. 우리는 본 세션에서 정가, 할인율, 리뷰 건수, 평점 등 다양한 변수들이 흥행 지표인 판매지수와 어떻게 얽혀 있는지 피어슨 상관계수 매트릭스를 기반으로 상세히 짚어볼 것입니다. 이를 통해 실질적으로 판매 성과를 자극하는 단 하나의 핵심 레버리지가 무엇인지 과학적으로 증명하겠습니다.",
        
        # 14. 변수 간 상관관계 매트릭스
        "지금 보시는 상관관계 매트릭스 히트맵 시각화 결과는 가격, 할인율, 리뷰 및 판매 실적 간의 수치적 조화를 명확하게 보여줍니다. 교차 영역 중 유독 붉게 표시되는 지점이 바로 리뷰 건수와 판매지수의 교차점입니다. 상관계수가 0.48로 강력하게 잡혀 양의 선형 관계를 입증합니다. 반면 할인율과 판매지수, 그리고 도서 정가와 판매지수 간의 교차점은 거의 무채색에 가까워 통계적인 선형 관계가 존재하지 않음을 실증적으로 대변합니다.",
        
        # 15. 상관계수 상세 분석 및 통찰
        "상관계수 테이블의 세부 데이터를 자세히 짚어드리겠습니다. 리뷰 건수와 판매지수 간의 계수는 0.48로 매우 유의미한 양의 상관성을 보입니다. 반면 도서 정가와 판매지수는 0.06으로 사실상 무상관에 가깝습니다. 이는 책의 가격이 비싸다고 해서 판매에 부정적인 영향을 미치지 않음을 설명합니다. 할인율 역시 0.05로 무상관입니다. 즉, 도서정가제 범위 내에서 10%의 기본 할인을 해주는 것은 독자들의 구매 결정에 결정적인 기여를 하지 못하며, 책의 학술적/실무적 가치가 훨씬 지배적인 구매 결정 요인임을 시사합니다.",
        
        # 16. 리뷰 건수와 판매지수의 선순환 메커니즘
        "리뷰 건수와 판매 성과 간의 '선순환 메커니즘'에 대해 자세히 논의하겠습니다. 책이 출간된 후 초기 마케팅이나 품질 보증을 통해 판매가 촉진되면 독자들의 자발적인 리뷰 서평이 차례로 적립됩니다. 축적된 리뷰의 볼륨은 신규 독자 유입 시에 강력한 사회적 증거 역할을 하여 구매 망설임을 지워줍니다. 실제로 분석 결과 리뷰 건수가 100건을 넘어서면 판매지수가 기하급수적으로 튀어 오르는 티핑 포인트 임계점이 확인됩니다. 따라서 초기 마케팅 리소스를 리뷰 100건 조기 달성에 조율하는 설계가 마케팅의 정답입니다.",
        
        # 17. 평점(리뷰총점) 데이터의 변별력 한계와 대안
        "리뷰의 만족도인 평점(리뷰총점) 데이터가 왜 판매지수와 상관계수가 0.16에 불과한지 실태를 알려드리겠습니다. 분석 결과, 베스트셀러 반열에 올라온 대부분의 IT 도서들은 평점이 9.5점에서 9.9점 사이에 빽빽하게 쏠려 있습니다. 중앙값이 무려 9.7점입니다. 즉, 평점 데이터 자체는 평점 상향 평준화로 인해 변별력 지표로서의 기능을 잃었습니다. 독자들 역시 평점 소수점 차이보다는 해당 평점이 누적된 '절대 리뷰 수'를 품질 신뢰의 직접적인 척도로 삼기 때문에 리뷰 개수 확보가 대안이 됩니다.",
        
        # 18. 도서 가격과 판매 성과의 상호 독립성
        "상관계수 0.06이 증명하는 도서 정가와 흥행의 상호 독립성 파트입니다. 많은 기획자들이 책 가격이 오르면 판매가 줄어들 것이라 우려하지만 데이터는 그것이 편견임을 입증합니다. IT 전문 독자들은 본인에게 유용한 코드가 들어있다면 35,000원 이상의 고가 도서도 주저 없이 기꺼이 구매합니다. 억지로 가격을 내리려 페이지나 품질을 깎기보다, 책의 실용성을 높여 고단가 고마진을 책정하는 것이 합리적인 전략입니다.",
        
        # 19. [간지] 출판사 영향력 및 트렌드 분석
        "네 번째 파트인 출판사 영향력 및 트렌드 분석 세션을 시작하겠습니다. 우리는 이번 장에서 IT 도서 시장을 점유하고 있는 메이저 10대 출판사들의 정량적 등록 도서 수와 권당 평균 판매 효율을 비교해 보고, 이들의 운영 전략(양적 다각화 vs 질적 메가 히트)의 차이점을 파악해 보겠습니다. 아울러 고객 평점 패턴 및 연도별 트렌드의 시계열 분포까지 심층 탐색하겠습니다.",
        
        # 20. 출판사 영향력 비교 시각화
        "화면의 차트는 출판사별 등록 도서 수와 평균 판매지수를 동시 매핑한 결과입니다. 상단 그래프에서는 한빛미디어가 132권으로 베스트셀러 목록을 독점하며 1위를 차지하고 있습니다. 그러나 하단 그래프인 평균 판매 효율로 가 보면, 이지스퍼블리싱(평균 약 1만 점)과 골든래빗(평균 약 6.4K)이 압도적인 피크로 올라선 것을 볼 수 있습니다. 이는 출판사 규모에 따라 시장을 침투하는 비즈니스 모델이 전혀 다르게 설계되어 있음을 대변합니다.",
        
        # 21. Top 10 출판사 주요 성과 상세
        "출판사별 세부 성과 테이블을 해설해 드리겠습니다. 한빛미디어는 132권으로 점유율 1위의 지배력을 유지하며 평균 지수 4,213으로 브랜드 파워를 지켰습니다. 한편 이지스퍼블리싱은 39권 등록에 권당 평균 9,931점이라는 압도적 1위의 질적 판매 효율을 보였고, 골든래빗 역시 43권에 6,368점으로 메가 히트 전략을 성공시켰습니다. 반면 자격증 수험서 위주의 영진닷컴(1,189)이나 대학교재 위주의 길벗캠퍼스(557) 등은 평균 판매 화력은 낮아 특정 목적성 타겟층에만 소구되고 있음을 알 수 있습니다.",
        
        # 22. 양적 점유율 vs 질적 성과의 비즈니스 시사점
        "출판사 성과 상세 데이터가 주는 교훈은 브랜드 다각화의 결론입니다. 한빛미디어는 다품종 소량 생산 혹은 넓은 카테고리 장악을 통한 롱테일 전방위 전략을 폅니다. 반면 이지스퍼블리싱이나 골든래빗은 독자 커뮤니티(Do it! 등) 및 표적 기획을 통해 런칭하는 도서마다 대성공을 거두는 메가 히트 전략에 충실합니다. 따라서 당사 원고 기획 시 대중적인 인프라 유통이 필요하다면 한빛미디어가, 초기 임팩트 중심 마케팅에는 이지스퍼블리싱이나 골든래빗 채널의 파트너십이 비즈니스적으로 현명합니다.",
        
        # 23. 고객 평점 및 리뷰 패턴 시각화
        "고객 평점과 리뷰 건수 및 판매지수의 버블 시각화 해설입니다. 고성능 성과를 거두는 도서군(판매지수 4K 이상)은 예외 없이 버블 색상이 오렌지 및 진한 빨간색(평점 9.5~10점) 사이에 강하게 집적되어 있습니다. 만족도가 우수한 도서들만 독자의 자발적 리뷰 바이럴이 붙어 판매가 지속될 수 있음을 증명합니다. 반면 품질이 검증되지 않은 책들은 누적 리뷰 수가 극도로 제한되고 바닥에 깔려 소외되므로, 집필 단계부터 정밀 감수를 통해 평점 9.0선 이하로 떨어지는 것을 철저히 사전 차단해야 합니다.",
        
        # 24. 연도별 출간 추이 및 트렌드 시각화
        "출간 연도에 따른 베스트셀러 진입 빈도 분석 결과입니다. 히스토그램에서 보시는 바와 같이 최근 2~3년(2024~2026년)에 런칭된 최신 신간 비중이 80% 이상으로 점유되어 있습니다. IT 기술은 트렌드 변화 주기가 극도로 짧아 신기술 버전 업그레이드에 따라 도서의 수명이 짧아지는 경향이 강하기 때문입니다. 단, 엑셀, C/파이썬 기초, OS/네트워크 등의 기초 학술서는 스테디셀러 생명력을 가지므로, 당사는 매년 신간 부스팅 전략과 기초 학술서 포트폴리오를 균형 있게 구성해야 매출 안정성이 확보됩니다.",
        
        # 25. [간지] 핵심 키워드 분석 및 비즈니스 추천 전략
        "다섯 번째 파트이자 마지막인 핵심 키워드 분석 및 비즈니스 추천 전략 세션을 시작하겠습니다. 우리는 텍스트 마이닝 기법인 TF-IDF 가중치를 적용하여, IT 서적 독자층의 실제 기술 수요 상위 30대 키워드를 순위별로 규명하고 이를 종합한 당사 차기 기획 전략 및 마케팅 채널 액션 플랜을 제시하겠습니다.",
        
        # 26. TF-IDF 기반 상위 1~15위 핵심 키워드 분석
        "TF-IDF 가중치 상위 1~15위 키워드 테이블 분석입니다. 가중치 1위는 0.421을 얻은 'AI'로, 현재 시장의 흐름이 인공지능 담론에 지배되고 있음을 보여줍니다. 이어 '파이썬(0.354)', '엑셀(0.312)', '코딩(0.287)' 등이 기초 스테디 영역을 장악하고 있습니다. 특징적인 점은 '클로드(0.254)', '제미나이(0.231)', '에이전트(0.210)', '바이브(0.125)' 등이 매우 높게 랭크되었다는 사실인데, 독자 수요가 AI를 이용해 즉각 풀스택 앱을 개발하는 생산성 코딩으로 완전히 이전되었음을 시사합니다.",
        
        # 27. TF-IDF 기반 상위 16~30위 차순위 키워드 분석
        "이어서 가중치 16위부터 30위까지의 키워드 맵입니다. '데이터(0.109)', '웹(0.098)', '자동화(0.091)' 등이 개발 실무 영역을 대변하고 있으며, '수업(0.076)'과 '에듀테크(0.038)' 키워드는 초중고 교사와 학부모들을 타겟으로 한 캔바, 노션 등 오피스/디자인 활용 도서군이 시장 내에서 대단히 고정적이고 강력한 특수 소비군 세그먼트를 탄탄하게 이루고 있음을 실증합니다. 이 키워드들을 적절히 조합하여 차기 서적의 기획 및 네이밍을 설계해야 합니다.",
        
        # 28. 키워드로 본 IT 트렌드 요약
        "키워드 분석을 통해 도출한 IT 도서 트렌드의 3대 요약입니다. 첫째, 인공지능(AI)의 전방위적 침투입니다. 제미나이 활용법, 클로드 코드 활용 등 도구 최적화 지식이 중심에 서 있습니다. 둘째, 업무 자동화 실무 니즈 팽창입니다. 직장인들이 엑셀이나 n8n 노코드 툴을 AI와 결합하여 생산성을 극대화하려는 경향이 강합니다. 셋째, 자연어를 매개로 한 바이브 코딩 패러다임의 확산입니다. AI 에이전트와 대화하며 풀스택 서비스를 조기 구축하는 실전 가이드가 인기를 끌고 있습니다.",
        
        # 29. 비즈니스 액션 플랜 - 기획/가격/마케팅/유통
        "최종 분석 결과를 바탕으로 도출한 4대 핵심 비즈니스 액션 플랜입니다. 첫째, 기획 부문에서는 단편적 챗봇을 넘어 시스템 관점의 'n8n 자동화' 및 '멀티 에이전트 설계'에 자원을 배분하십시오. 둘째, 가격 부문에서는 마진 보존을 위해 **24,000~28,000원**의 고품질 표준 단가를 유지하십시오. 셋째, 마케팅 부문에서는 리뷰 건수 상관성(0.48)을 고려하여 출간 초기 1개월 내 **리뷰 50~100건 조기 돌파**에 집중 투자하십시오. 넷째, 유통은 양적 커버리지와 질적 메가 히트에 맞춰 한빛 혹은 이지스/골든래빗 채널을 차별화 타겟팅하십시오.",
        
        # 30. 종합 결론 및 마무리
        "이상으로 YES24 IT/모바일 베스트셀러 데이터 분석 결과를 모두 마치겠습니다. 우리는 데이터를 통해 IT 서적 시장의 파레토 롱테일 구조, 실무 서적의 적정 가격 및 마진 안정성, 리뷰 개수가 판매에 미치는 선순환 효과, 메이저 출판사별 포지셔닝 차이점, AI 전방위 지배에 따른 바이브 코딩의 트렌드 부상을 실증적으로 확인했습니다. 본 분석 모델과 제안해 드린 4대 액션 플랜이 당사 출판/마케팅 실무에 실질적인 매출 성장과 성공률 제고의 나침반이 되길 기원합니다. 경청해 주셔서 대단히 감사합니다."
    ]
    
    # 디렉토리 확인
    md_dir = os.path.join("yes24", "docs")
    
    # ----------------------------------------------------
    # 슬라이드 생성 루프
    # ----------------------------------------------------
    
    # Slide 1: 표지
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_LEAD_BG)
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = "YES24 IT/모바일 베스트셀러 데이터 분석"
    run1.font.name = 'Arial'
    run1.font.size = Pt(42)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_LEAD_TEXT
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(14)
    run2 = p2.add_run()
    run2.text = "탐색적 데이터 분석(EDA) 및 비즈니스 인사이트 도출"
    run2.font.name = '나눔고딕'
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0xD9, 0xE1, 0xF2)
    
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)
    run3 = p3.add_run()
    run3.text = "발표자: 데이터 분석 연구원  |  대상 데이터: IT/모바일 분야 도서 869건"
    run3.font.name = '나눔고딕'
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
    
    # Slide 2: 목차
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "목차 (Table of Contents)")
    # 좌측 목차 타이틀 카드
    add_card(slide, 0.8, 1.8, 4.0, 4.5, COLOR_CARD)
    add_text_box(slide, 1.0, 2.8, 3.6, 2.0, [
        {"text": "발표 목차 안내\n\n", "font": "Arial", "size": 24, "bold": True, "color": COLOR_TITLE},
        {"text": "본 분석 모델은 8개 핵심 요인에 맞추어 유기적으로 설계되었습니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
    ])
    # 우측 카드 배치 (목차 리스트)
    add_card(slide, 5.2, 1.8, 7.3, 4.5, COLOR_WHITE)
    toc_items = [
        "분석 개요 및 기초 통계",
        "도서 가격대 분포 분석",
        "변수 간 상관관계 분석",
        "출판사 영향력 및 성과 분석",
        "고객 리뷰 및 평점 패턴 분석",
        "연도별 출간 추이 및 라이프사이클",
        "TF-IDF 기반 핵심 키워드 도출",
        "비즈니스 추천 전략 및 결론"
    ]
    for i, item in enumerate(toc_items):
        y_pos = 2.1 + i * 0.5
        add_flat_icon(slide, 5.5, y_pos, 0.35, "diamond")
        add_text_box(slide, 6.0, y_pos - 0.05, 6.0, 0.4, [
            {"text": item, "font": "나눔고딕", "size": 15, "bold": True, "color": COLOR_TEXT}
        ])
    
    # Slide 3: [간지 1] 개요 및 기술통계
    create_lead_slide("분석 개요 및 기초 통계", "세션 1: 데이터 소개 및 기초 통계량 파악")
    
    # Slide 4: 분석 배경 및 데이터 소개
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "분석 배경 및 데이터 소개")
    # 인포그래픽형 2열 구조
    add_card(slide, 0.8, 1.8, 5.5, 4.5, COLOR_CARD)
    add_flat_icon(slide, 1.1, 2.0, 0.4, "diamond")
    add_text_box(slide, 1.0, 2.5, 5.1, 3.6, [
        {"text": "데이터 수집 개요\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 수집 대상: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "YES24 IT/모바일 베스트셀러 리스트\n", "font": "나눔고딕", "size": 14},
        {"text": "• 데이터 규모: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "총 869건의 독자적인 도서 레코드\n", "font": "나눔고딕", "size": 14},
        {"text": "• 수집 변수: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "정가, 판매가, 판매지수, 리뷰 등 15개 변수\n\n", "font": "나눔고딕", "size": 14},
        {"text": "빠르게 변화하는 IT 기술 시장 트렌드에 따라 IT 서적 소비자들의 실질적 니즈를 정량적으로 파악하기 위한 기초 데이터셋을 마련했습니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT, "italic": True}
    ])
    add_card(slide, 7.0, 1.8, 5.5, 4.5, COLOR_WHITE)
    add_flat_icon(slide, 7.3, 2.0, 0.4, "star")
    add_text_box(slide, 7.2, 2.5, 5.1, 3.6, [
        {"text": "분석 핵심 목적\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 시장 지배력 확인: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "메이저 출판사들의 실제 유통 성과 분석\n", "font": "나눔고딕", "size": 14},
        {"text": "• 흥행 메커니즘 규명: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "판매 성과와 리뷰 패턴 간 연관성 증명\n", "font": "나눔고딕", "size": 14},
        {"text": "• 마케팅 액션 도출: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "신간 출간 시 이익 극대화를 위한 가격 정책 수립 및 예산 배분 가이드 확보", "font": "나눔고딕", "size": 14}
    ])
    
    # Slide 5: 데이터 전처리 및 정제
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "데이터 전처리 및 정제 요약")
    # 가로 프로세스형 4단 배치
    processes = [
        {"title": "1. 저자 결측 정제", "desc": "저자 누락 2건에 대해\n'저자 미상' 대체\n통계 무결성 유지"},
        {"title": "2. 포인트 수치 변환", "desc": "'원' 기호 및 콤마 제거\n정수형 데이터 변환\n포인트 평균 연산 기반 마련"},
        {"title": "3. 할인율 소수 매핑", "desc": "10% 단위 할인 수치를\n소수(예: 0.10)로 변환\n엑셀 동적 호환성 확보"},
        {"title": "4. 무결성 검증", "desc": "가상환경 `.venv`에서\n중복 검사 실시\n중복 레코드 0건 확인"}
    ]
    for i, proc in enumerate(processes):
        x = 0.8 + i * 2.95
        add_card(slide, x, 2.0, 2.8, 4.0, COLOR_CARD)
        add_flat_icon(slide, x + 0.2, 2.2, 0.35, "diamond")
        add_text_box(slide, x + 0.1, 2.7, 2.6, 3.1, [
            {"text": proc["title"] + "\n\n", "font": "Arial", "size": 16, "bold": True, "color": COLOR_TITLE},
            {"text": proc["desc"], "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
        ])
        
    # Slide 6: 거시적 주요 요약 지표 (KPI)
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "거시적 주요 요약 지표 (KPI)")
    # 대형 숫자 카드 3개 배치
    kpis = [
        {"title": "총 베스트셀러 도서 수", "val": "869", "unit": " 건", "bg": COLOR_CARD, "accent": COLOR_TITLE, "icon": "diamond"},
        {"title": "도서 평균 실판매가", "val": "23,395", "unit": " 원", "bg": COLOR_WHITE, "accent": COLOR_ACCENT, "icon": "star"},
        {"title": "누적 판매지수 총합", "val": "2,347,452", "unit": " 점", "bg": COLOR_CARD, "accent": COLOR_TITLE, "icon": "cloud"}
    ]
    for i, kpi in enumerate(kpis):
        x = 0.8 + i * 3.95
        add_card(slide, x, 2.0, 3.8, 3.8, kpi["bg"])
        # 데코레이션 사각형
        dec = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.0), Inches(3.8), Inches(0.15))
        dec.fill.solid()
        dec.fill.fore_color.rgb = kpi["accent"]
        dec.line.fill.background()
        
        # 플랫 아이콘
        add_flat_icon(slide, x + 1.7, 2.3, 0.4, kpi["icon"])
        
        add_text_box(slide, x + 0.2, 2.8, 3.4, 2.7, [
            {"text": kpi["title"] + "\n\n", "font": "나눔고딕", "size": 15, "color": COLOR_TEXT},
            {"text": kpi["val"], "font": "Arial", "size": 34, "bold": True, "color": COLOR_TITLE},
            {"text": kpi["unit"], "font": "나눔고딕", "size": 18, "bold": True, "color": COLOR_TEXT}
        ], align=PP_ALIGN.CENTER)
        
    # Slide 7: 수치 데이터 기술통계 요약
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "수치형 데이터 기술통계 요약")
    # 표 배치
    table_shape = slide.shapes.add_table(7, 8, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.2))
    table = table_shape.table
    # 열 너비 설정
    table.columns[0].width = Inches(1.5)
    for col_i in range(1, 8):
        table.columns[col_i].width = Inches(1.46)
        
    headers = ["통계량", "할인율", "판매가(원)", "정가(원)", "포인트(원)", "판매지수", "리뷰건수", "리뷰평점"]
    for col_i, h_text in enumerate(headers):
        cell = table.cell(0, col_i)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = COLOR_WHITE
                
    row_data = [
        ["평균", "8.4%", "23,395", "25,560", "1,135", "2,701", "18.6", "7.55"],
        ["최솟값", "0.0%", "5,625", "6,250", "0", "60", "0", "0.0"],
        ["중앙값", "10.0%", "22,500", "24,400", "1,100", "1,062", "9", "9.7"],
        ["최댓값", "10.0%", "59,400", "66,000", "3,300", "79,368", "388", "10.0"],
        ["합계", "-", "20,330,490", "22,211,800", "986,510", "2,347,452", "16,168", "6,564.9"],
        ["데이터 수", "869", "869", "869", "869", "869", "869", "869"]
    ]
    for row_i, row_val in enumerate(row_data, 1):
        for col_i, val in enumerate(row_val):
            cell = table.cell(row_i, col_i)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD if row_i % 2 == 0 else COLOR_WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_i == 0 else PP_ALIGN.RIGHT
                for r in p.runs:
                    r.font.name = "나눔고딕"
                    r.font.size = Pt(11)
                    r.font.color.rgb = COLOR_TEXT
                    
    # Slide 8: 기술통계에서 발견된 비즈니스 롱테일
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "기술통계에서 발견된 비즈니스 롱테일")
    # 인포그래픽형 사각형 (파레토 법칙 시각화)
    add_card(slide, 0.8, 1.8, 6.0, 4.5, COLOR_CARD)
    add_flat_icon(slide, 1.1, 2.0, 0.4, "cloud")
    add_text_box(slide, 1.1, 2.5, 5.4, 3.5, [
        {"text": "평균과 중앙값의 엄청난 격차\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 판매지수 평균: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "2,701 점\n", "font": "나눔고딕", "size": 14},
        {"text": "• 판매지수 중앙값: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "1,062 점 (평균 대비 2.5배 격차)\n", "font": "나눔고딕", "size": 14},
        {"text": "• 리뷰건수 평균: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "18.6 건\n", "font": "나눔고딕", "size": 14},
        {"text": "• 리뷰건수 중앙값: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "9 건 (평균 대비 2배 격차)\n\n", "font": "나눔고딕", "size": 14},
        {"text": "이 통계적 불일치는 베스트셀러 내에서도 극히 소수의 '메가 히트작'이 전체 성과를 이끌어 가고 있음을 증명합니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT, "italic": True}
    ])
    
    # 우측 파레토 80:20 레이아웃 박스
    add_card(slide, 7.2, 1.8, 5.3, 4.5, COLOR_TITLE)
    add_flat_icon(slide, 7.5, 2.0, 0.4, "star")
    add_text_box(slide, 7.5, 2.5, 4.7, 3.5, [
        {"text": "IT 서적 시장의 파레토 법칙\n\n", "font": "Arial", "size": 22, "bold": True, "color": COLOR_BG},
        {"text": "상위 20% 도서가 누적 판매지수의 80% 이상을 차지하는 전형적인 롱테일 구조를 보입니다. 마케팅 비용을 모든 라인업에 고르게 분산하기보다 킬러 타이틀의 조기 부스팅에 집중 투자하는 채널 집중화가 필수적입니다.", "font": "나눔고딕", "size": 14, "color": COLOR_BG}
    ])
    
    # Slide 9: [간지 2] 도서 가격대 분포 분석
    create_lead_slide("도서 가격대 분포 분석", "세션 2: IT 도서 가격의 분포 범위와 가격 정책 제언")
    
    # Slide 10: 도서 가격 분포 시각화
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_image_slide_content(slide, "price_dist.png", md_dir, "도서 가격대 분포 분석 (시각화)", [
        {"text": "가격 집중 구간 분석\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 주요 정가대: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "15,000원 ~ 30,000원 구간\n", "font": "나눔고딕", "size": 14},
        {"text": "• 점유율: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "베스트셀러 전체 도서의 70% 이상\n", "font": "나눔고딕", "size": 14},
        {"text": "• 고단가 꼬리 분포: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "35,000원 이상의 전문 서적도 넓게 포진\n\n", "font": "나눔고딕", "size": 14},
        {"text": "분석 결과 IT 서적은 일반 대중서보다 단가가 높게 포지셔닝되는 경향이 있습니다. 이는 인쇄 볼륨과 소스 코드 예제 삽입 등 기술 도서 제작 과정의 단가가 단가 분포에 고스란히 반영된 흔적입니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TITLE}
    ])
    
    # Slide 11: 가격대 구간별 분포의 비즈니스적 해석
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "가격대 구간별 분포의 비즈니스 해석")
    # 3단 세로 컬럼 카드
    cats = [
        {"title": "저가 구간 (1.5만 이하)", "desc": "• 대상: 얇은 모바일 매뉴얼, 자격증 기출 수험서, 전자책\n• 특징: 가볍게 읽는 용도 위주로 단가가 1.5만 원 미만으로 강제 차단됨", "icon": "diamond"},
        {"title": "표준 구간 (1.5만 ~ 3만)", "desc": "• 대상: 프로그래밍 언어 입문, 프레임워크 실습서, 오피스 책\n• 특징: 베스트셀러 대다수가 밀집한 대중 핵심 캐시카우 구간", "icon": "star"},
        {"title": "고가 구간 (3.5만 이상)", "desc": "• 대상: AI 대학 교과서, 인프라 아키텍처 원서, 전문 전공서\n• 특징: 높은 난이도와 방대한 페이지 볼륨으로 3.5만 원 이상의 가격 포지셔닝", "icon": "arrow"}
    ]
    for i, cat in enumerate(cats):
        x = 0.8 + i * 3.95
        add_card(slide, x, 2.0, 3.8, 4.5, COLOR_CARD)
        add_flat_icon(slide, x + 0.3, 2.2, 0.35, cat["icon"])
        add_text_box(slide, x + 0.2, 2.7, 3.4, 3.6, [
            {"text": cat["title"] + "\n\n", "font": "Arial", "size": 16, "bold": True, "color": COLOR_TITLE},
            {"text": cat["desc"], "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
        ])
        
    # Slide 12: 최적 가격 포지셔닝 및 기획 전략
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "최적 가격 포지셔닝 및 기획 제언")
    # 체크박스 & 추천 태그 레이아웃
    add_card(slide, 0.8, 1.8, 5.5, 4.5, COLOR_CARD)
    add_flat_icon(slide, 1.1, 2.0, 0.4, "smile")
    add_text_box(slide, 1.0, 2.5, 5.1, 3.5, [
        {"text": "고객의 가격 감수성 분석\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 고관여 타겟의 특징: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "IT 소비자는 단순 가격 할인보다 서적이 전달하는 기술적 가치와 지식의 유용성에 따라 도서를 구매합니다.\n\n", "font": "나눔고딕", "size": 13},
        {"text": "• 가격과 이익률 균형: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "단가 인하를 위해 내용이나 코드를 빼는 편집을 하기보다, 확실하게 품질을 보강해 적정 가격을 받는 가치 마케팅이 비즈니스적으로 성공률이 높습니다.", "font": "나눔고딕", "size": 13}
    ])
    add_card(slide, 7.0, 1.8, 5.5, 4.5, COLOR_TITLE)
    add_flat_icon(slide, 7.3, 2.0, 0.4, "arrow")
    add_text_box(slide, 7.2, 2.5, 5.1, 3.5, [
        {"text": "차기 도서 가격 가이드\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_BG},
        {"text": "• 신간 입문/실무 서적 권장 정가: ", "font": "나눔고딕", "size": 14, "bold": True, "color": COLOR_BG},
        {"text": "24,000원 ~ 28,000원 대\n\n", "font": "Arial", "size": 22, "bold": True, "color": COLOR_ACCENT},
        {"text": "이 범위는 제조 마진을 안전하게 방어하고, 독자의 심리적 저항선을 최소화하는 최적의 단가 타협점입니다. 고급 예제와 가독성을 보장하여 가격 설득력을 높여야 합니다.", "font": "나눔고딕", "size": 13, "color": COLOR_BG}
    ])
    
    # Slide 13: [간지 3] 변수 간 상관관계 분석
    create_lead_slide("변수 간 상관관계 분석", "세션 3: 흥행에 영향을 미치는 핵심 요인들의 선형 상관성 규명")
    
    # Slide 14: 변수 간 상관관계 매트릭스
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_image_slide_content(slide, "correlation_matrix.png", md_dir, "변수 간 상관관계 매트릭스 (시각화)", [
        {"text": "상관관계 주요 발견\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 리뷰건수 & 판매지수: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "상관계수 0.48 (가장 유의미한 양의 상관성)\n", "font": "나눔고딕", "size": 14},
        {"text": "• 정가 & 판매지수: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "상관계수 0.06 (상관관계 없음)\n", "font": "나눔고딕", "size": 14},
        {"text": "• 할인율 & 판매지수: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "상관계수 0.05 (상관관계 없음)\n\n", "font": "나눔고딕", "size": 14},
        {"text": "이 통계는 '싸고 할인을 많이 한다고 해서 책이 더 흥행하지 않는다'는 점과 '독자의 리뷰 수'가 판매의 원동력임을 증명합니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TITLE}
    ])
    
    # Slide 15: 상관계수 상세 분석 및 통찰
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "변수 간 상관관계 상세 요약")
    # 테이블 배치
    table_shape = slide.shapes.add_table(5, 4, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.2))
    table = table_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(4.7)
    
    headers = ["상관 분석 변수 쌍", "상관계수", "관계 수준", "비즈니스 통찰 및 해석"]
    for col_i, h in enumerate(headers):
        cell = table.cell(0, col_i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(12)
                r.font.bold = True
                r.font.color.rgb = COLOR_WHITE
                
    row_data = [
        ["리뷰건수 & 판매지수", "0.48", "뚜렷한 양의 상관", "리뷰 볼륨 축적이 신규 독자 유입의 선순환 견인"],
        ["정가 & 판매지수", "0.06", "상관관계 없음", "도서 가격 높낮이는 베스트셀러 등극에 무영향"],
        ["리뷰평점 & 판매지수", "0.16", "매우 약한 상관", "평점 수치 차이보다 절대 서평 참여 개수가 더 중요"],
        ["할인율 & 판매지수", "0.05", "상관관계 없음", "도서정가제 내 10% 한도 할인은 차별성 제공 못함"]
    ]
    for row_i, row in enumerate(row_data, 1):
        for col_i, val in enumerate(row):
            cell = table.cell(row_i, col_i)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD if row_i % 2 == 0 else COLOR_WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_i in [0, 1, 2] else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "나눔고딕"
                    r.font.size = Pt(12)
                    r.font.color.rgb = COLOR_TEXT
                    
    # Slide 16: 리뷰 건수와 판매지수의 선순환 메커니즘
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "리뷰 누적과 판매지수의 선순환 루프")
    # 화살표 순환 구조를 나타내는 3단 가로 배치
    stages = [
        {"title": "1. 독자 구매 촉진", "desc": "초기 광고 및 서평단 활성화로 도서가 판매되기 시작함"},
        {"title": "2. 자발적 후기 누적", "desc": "독자 만족도를 보증하는 구매 리뷰와 서평 개수가 차례로 적립됨"},
        {"title": "3. 신뢰 상승 & 판매 재자극", "desc": "누적 리뷰가 신규 유입자들의 구매 장벽을 낮춰 판매량이 다시 튐"}
    ]
    for i, stage in enumerate(stages):
        x = 0.8 + i * 3.95
        add_card(slide, x, 2.0, 3.8, 4.5, COLOR_CARD)
        add_flat_icon(slide, x + 0.3, 2.2, 0.35, "arrow")
        add_text_box(slide, x + 0.2, 2.7, 3.4, 3.6, [
            {"text": stage["title"] + "\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
            {"text": stage["desc"], "font": "나눔고딕", "size": 14, "color": COLOR_TEXT}
        ])
        if i < 2:
            # 화살표 도형 추가
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 3.8), Inches(4.0), Inches(0.15), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_ACCENT
            arrow.line.fill.background()
            
    # Slide 17: 평점(리뷰총점) 데이터의 변별력 한계와 대안
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "평점 데이터의 변별력 한계와 대안")
    # 경고/주의 카드 디자인
    add_card(slide, 0.8, 1.8, 5.5, 4.5, COLOR_CARD)
    add_flat_icon(slide, 1.1, 2.0, 0.4, "diamond")
    add_text_box(slide, 1.0, 2.5, 5.1, 3.5, [
        {"text": "평점 상향 평준화 현상\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 평점 중앙값: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "9.7 점 (10점 만점)\n", "font": "나눔고딕", "size": 14},
        {"text": "• 평점 밀집도: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "대다수 활성 도서의 평점이 9.5~9.9점 사이에 존재\n\n", "font": "나눔고딕", "size": 14},
        {"text": "이 상향 평준화 분포로 인해 평점 수치 자체의 변별력(상관계수 0.16)은 매우 떨어집니다. 평점 소수점의 미미한 점수차는 독자의 구매 자극 요소가 되지 못합니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
    ])
    add_card(slide, 7.0, 1.8, 5.5, 4.5, COLOR_TITLE)
    add_flat_icon(slide, 7.3, 2.0, 0.4, "star")
    add_text_box(slide, 7.2, 2.5, 5.1, 3.5, [
        {"text": "마케팅적 대안\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_BG},
        {"text": "• 리뷰 개수(볼륨) 확보 집중: ", "font": "나눔고딕", "size": 14, "bold": True, "color": COLOR_BG},
        {"text": "독자는 몇 점인지보다 '얼마나 많은 사람들이 이 책을 읽고 호평했는가' 즉 리뷰의 양적 수치를 품질 신뢰의 직접 척도로 삼습니다. 따라서 마케팅 역량을 리뷰 건수 증대에 맞춰야 합니다.", "font": "나눔고딕", "size": 13, "color": COLOR_BG}
    ])
    
    # Slide 18: 도서 가격과 판매 성과의 상호 독립성
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "도서 정가와 판매지수의 독립성")
    # 대조형 2열 레이아웃
    add_card(slide, 0.8, 1.8, 5.5, 4.5, COLOR_WHITE)
    add_flat_icon(slide, 1.1, 2.0, 0.4, "cloud")
    add_text_box(slide, 1.0, 2.5, 5.1, 3.5, [
        {"text": "일반적 편견\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 가격 저항감에 대한 우려: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "도서 가격이 비쌀수록 잠재적 수요가 줄어 판매에 부정적 영향을 미칠 것이라는 막연한 예상이 흔히 존재합니다.\n\n", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT},
        {"text": "• 실판매가 인하 집착: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "이로 인해 편집 단가를 무리하게 낮춰 퀄리티를 타협하려는 유혹이 흔히 발생합니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
    ])
    add_card(slide, 7.0, 1.8, 5.5, 4.5, COLOR_CARD)
    add_flat_icon(slide, 7.3, 2.0, 0.4, "diamond")
    add_text_box(slide, 7.2, 2.5, 5.1, 3.5, [
        {"text": "실제 데이터의 진실\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 정가 & 판매지수 상관계수: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "0.06 (무작위 무상관)\n\n", "font": "Arial", "size": 22, "bold": True, "color": COLOR_TITLE},
        {"text": "정가와 흥행 실적 간의 관계는 전혀 관찰되지 않습니다. IT 전문 도서 독자들은 지식의 실용적 활용 가치가 있다면 35,000원 이상의 고가 도서도 주저 없이 지출하기 때문에 고마진 정책이 유효합니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
    ])
    
    # Slide 19: [간지 4] 출판사 영향력 및 트렌드 분석
    create_lead_slide("출판사 영향력 및 트렌드 분석", "세션 4: 메이저 출판사들의 시장 점유율과 운영 효율성 비교")
    
    # Slide 20: 출판사 영향력 비교 시각화
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_image_slide_content(slide, "top_publishers.png", md_dir, "출판사별 점유율 및 판매지수 (시각화)", [
        {"text": "출판사별 성과 비교\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 최다 베스트셀러 등록: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "한빛미디어 (132권, 독보적 1위)\n", "font": "나눔고딕", "size": 14},
        {"text": "• 최고 판매 효율 출판사: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "이지스퍼블리싱 (평균 지수 약 1만 점)\n", "font": "나눔고딕", "size": 14},
        {"text": "• 강소 히트 출판사: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "골든래빗 (평균 지수 6,368점)\n\n", "font": "나눔고딕", "size": 14},
        {"text": "시장의 총 도서 보유 비중과 권당 개별 판매 성과는 다르게 나타납니다. 출판사 규모에 따라 시장을 포지셔닝하는 방식의 극명한 전략적 대조를 보입니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TITLE}
    ])
    
    # Slide 21: Top 10 출판사 주요 성과 상세
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "Top 10 출판사별 세부 성과")
    # 테이블 배치
    table_shape = slide.shapes.add_table(11, 4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(6.2)
    
    headers = ["출판사명", "등록 도서 수", "평균 판매지수", "비즈니스 영향력 및 세부 성격"]
    for col_i, h in enumerate(headers):
        cell = table.cell(0, col_i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.color.rgb = COLOR_WHITE
                
    row_data = [
        ["한빛미디어", "132", "4,213", "베스트셀러 보유량 1위의 명실상부 지배적인 IT 메이저 브랜드"],
        ["골든래빗", "43", "6,368", "최신 인공지능/에이전틱 코딩 분야 신간 메가 히트로 고효율 달성"],
        ["이지스퍼블리싱", "39", "9,931", "권당 평균 지수 1위, 독보적인 팬덤과 프랜차이즈 파워 보유"],
        ["길벗", "37", "2,752", "인프라, 기초 핵심 기술 교과서 라인업의 안정적인 장기 유통"],
        ["제이펍", "31", "5,160", "실무 엑셀 및 실용 컴퓨터 서적 등 뛰어난 판매지수 효율 증명"],
        ["위키북스", "30", "3,115", "트렌드 반영 신간을 적기에 출시하여 시장 점유 방어"],
        ["앤써북", "22", "4,524", "교사 대상 에듀테크 및 오피스 영역의 독자층 집중 공략"],
        ["영진닷컴", "17", "1,189", "전통 수험서 및 기초 활용 영역 위주로 판매지수 편차 낮음"],
        ["정보문화사", "16", "1,029", "전통 클래식 컴퓨터 서적으로 무난한 판매 흐름 유지"],
        ["길벗캠퍼스", "14", "557", "학기별 교재/학술 중심 출판으로 대중적 성과는 다소 한계"]
    ]
    for row_i, row in enumerate(row_data, 1):
        for col_i, val in enumerate(row):
            cell = table.cell(row_i, col_i)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD if row_i % 2 == 0 else COLOR_WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_i in [1, 2] else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "나눔고딕"
                    r.font.size = Pt(10)
                    r.font.color.rgb = COLOR_TEXT
                    
    # Slide 22: 양적 점유율 vs 질적 성과의 비즈니스 시사점
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "양적 점유율 vs 질적 성과의 대조")
    # 비교 평형 대칭 카드
    add_card(slide, 0.8, 1.8, 5.5, 4.5, COLOR_CARD)
    add_flat_icon(slide, 1.1, 2.0, 0.4, "diamond")
    add_text_box(slide, 1.0, 2.5, 5.1, 3.5, [
        {"text": "양적 다각화 브랜드 모델\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 대표 주자: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "한빛미디어\n", "font": "나눔고딕", "size": 14},
        {"text": "• 출판 특징: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "IT 전반의 촘촘한 라인업 설계를 통해 시장의 모든 니즈를 장악하는 롱테일 전방위 유통망 모델\n\n", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT},
        {"text": "• 권장 타겟: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "광범위한 대학 교재, 기관 및 공공 납품, 대형 유통 커버리지가 최우선인 기획물에 적합", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
    ])
    add_card(slide, 7.0, 1.8, 5.5, 4.5, COLOR_WHITE)
    add_flat_icon(slide, 7.3, 2.0, 0.4, "star")
    add_text_box(slide, 7.2, 2.5, 5.1, 3.5, [
        {"text": "질적 메가 히트 브랜드 모델\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 대표 주자: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "이지스퍼블리싱, 골든래빗\n", "font": "나눔고딕", "size": 14},
        {"text": "• 출판 특징: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "도서 1권당 폭발적인 판매 집중력을 유도하는 타겟 커뮤니티 팬덤형 모델\n\n", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT},
        {"text": "• 권장 타겟: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "초기 바이럴과 독자 소통이 극단적으로 중요한 핵심 AI 트렌드 및 최신 기술 신용서에 적합", "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
    ])
    
    # Slide 23: 고객 평점 및 리뷰 패턴 시각화
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_image_slide_content(slide, "sales_vs_reviews.png", md_dir, "고객 리뷰 분포와 판매지수 (시각화)", [
        {"text": "독자 만족도와 판매 화력\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 고성과 흥행작 특징: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "판매 성과가 높을수록 평점이 예외 없이 9.5~10점 사이에 초고밀도 집적\n\n", "font": "나눔고딕", "size": 13},
        {"text": "• 독자 불만 도서의 도태: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "만족도가 낮은 도서는 리뷰가 축적되지 못하고 판매지수 역시 즉시 바닥으로 수렴\n\n", "font": "나눔고딕", "size": 13},
        {"text": "결국 장기 베스트셀러 진입을 위해서는 집필 및 감수 감수 단계의 정밀한 품질 통제가 가장 본질적이고 선행되어야 할 가치임을 증명합니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TITLE}
    ])
    
    # Slide 24: 연도별 출간 추이 및 트렌드 시각화
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_image_slide_content(slide, "trend_by_year.png", md_dir, "연도별 출간 분포와 트렌드 (시각화)", [
        {"text": "IT 도서의 극히 짧은 생명 주기\n\n", "font": "Arial", "size": 18, "bold": True, "color": COLOR_TITLE},
        {"text": "• 신간 쏠림 현상: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "최근 2~3년(2024~2026년) 내 출간 도서가 전체의 80% 이상 장악\n\n", "font": "나눔고딕", "size": 13},
        {"text": "• 기술 트렌드 수명 주기: ", "font": "나눔고딕", "size": 14, "bold": True},
        {"text": "기술 스펙 업그레이드에 민감하여 구버전 실용서는 시장에서 즉각 도태\n\n", "font": "나눔고딕", "size": 13},
        {"text": "단, 엑셀, C/파이썬 기초, OS/네트워크 등의 클래식 기초서는 긴 생명력을 지니므로 안정적인 포트폴리오 믹스가 필수적입니다.", "font": "나눔고딕", "size": 13, "color": COLOR_TITLE}
    ])
    
    # Slide 25: [간지 5] 핵심 키워드 분석 및 비즈니스 추천 전략
    create_lead_slide("핵심 키워드 분석 및 비즈니스 추천 전략", "세션 5: 데이터 기반 IT 트렌드 도출 및 4대 실행 방안 제안")
    
    # Slide 26: TF-IDF 기반 상위 1~15위 핵심 키워드 분석
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "상위 1~15위 핵심 IT 키워드 분석")
    # 테이블 배치
    table_shape = slide.shapes.add_table(6, 4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(5.7)
    
    headers = ["순위", "핵심 키워드", "TF-IDF 가중치", "베스트셀러 시장 내 도서 성격 및 통찰"]
    for col_i, h in enumerate(headers):
        cell = table.cell(0, col_i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.color.rgb = COLOR_WHITE
                
    row_data = [
        ["1위 ~ 3위", "AI / 파이썬 / 엑셀", "0.421 / 0.354 / 0.312", "인공지능 실무 활용, 직장인 오피스 생산성이 주류 매출원 장악"],
        ["4위 ~ 6위", "코딩 / 클로드 / 제미나이", "0.287 / 0.254 / 0.231", "AI 에이전트 도구를 활용한 차세대 실무 코딩 테마 강세"],
        ["7위 ~ 9위", "에이전트 / 실무 / 활용법", "0.210 / 0.198 / 0.187", "단순 교양서보다 직접 결과를 내는 실용적인 프롬프트북 중심"],
        ["10위 ~ 12위", "노트북LM / 입문 / 자바", "0.165 / 0.154 / 0.142", "비전공자 타겟 IT 진입서 및 실무 백엔드 개발서의 공존"],
        ["13위 ~ 15위", "인프라 / 바이브 / 챗GPT", "0.131 / 0.125 / 0.118", "클라우드 서비스 구축 및 자연어 매개 코딩의 트렌드 부상"]
    ]
    for row_i, row in enumerate(row_data, 1):
        for col_i, val in enumerate(row):
            cell = table.cell(row_i, col_i)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD if row_i % 2 == 0 else COLOR_WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_i in [0, 1, 2] else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "나눔고딕"
                    r.font.size = Pt(11)
                    r.font.color.rgb = COLOR_TEXT
                    
    # Slide 27: TF-IDF 기반 상위 16~30위 차순위 키워드 분석
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "상위 16~30위 차순위 IT 키워드 분석")
    # 테이블 배치
    table_shape = slide.shapes.add_table(6, 4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(5.7)
    
    for col_i, h in enumerate(headers):
        cell = table.cell(0, col_i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(11)
                r.font.bold = True
                r.font.color.rgb = COLOR_WHITE
                
    row_data = [
        ["16위 ~ 18위", "데이터 / 웹 / 자동화", "0.109 / 0.098 / 0.091", "n8n 등 노코드 툴을 융합한 업무 자동화 설계 수요 증명"],
        ["19위 ~ 21위", "자격증 / 설계 / 수업", "0.087 / 0.082 / 0.076", "정보처리기사 컴활 시험서 및 교사용 에듀테크 캔바 수업 테마"],
        ["22위 ~ 24위", "모바일 / 딥러닝 / 네트워크", "0.071 / 0.068 / 0.063", "전통 전공/학습 기초 과목들의 장기 스테디 셀러화 작동"],
        ["25위 ~ 27위", "알고리즘 / 클라우드 / SQL", "0.059 / 0.054 / 0.049", "코딩테스트 대비 서적 및 데이터베이스 쿼리 분석 수요"],
        ["28위 ~ 30위", "캔바 / 보안 / 에듀테크", "0.045 / 0.041 / 0.038", "초중고 교육용 및 디자인 콘텐츠 제작 툴의 탄탄한 시장 구축"]
    ]
    for row_i, row in enumerate(row_data, 1):
        for col_i, val in enumerate(row):
            cell = table.cell(row_i, col_i)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_CARD if row_i % 2 == 0 else COLOR_WHITE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if col_i in [0, 1, 2] else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.name = "나눔고딕"
                    r.font.size = Pt(11)
                    r.font.color.rgb = COLOR_TEXT
                    
    # Slide 28: 키워드로 본 IT 트렌드 요약
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "핵심 키워드로 본 IT 소비 트렌드")
    # 3단 리본형 정보 태그
    trends = [
        {"title": "인공지능(AI)의 전방위 침투", "desc": "• AI, 클로드, 제미나이 등 최신 생성형 AI 활용서가 1위 지배\n• 단순 질의를 넘어 실제 에이전트를 조율하는 깊이 있는 지식이 핵심 트렌드", "icon": "star"},
        {"title": "실무 자동화와 오피스 결합", "desc": "• 엑셀, 자동화, n8n, 캔바 키워드의 강세\n• 비전공자 및 일반 직장인들이 AI를 도구로 사용하려는 강력한 실용성 수요 존재", "icon": "cloud"},
        {"title": "바이브 코딩 패러다임 부상", "desc": "• 클로드 코드로 빠르게 코딩을 런칭하는 실습서 선호\n• 복잡한 코드 문법 교육보다 개발 결과물을 신속하게 빌드하는 실전형 가이드 부상", "icon": "arrow"}
    ]
    for i, trend in enumerate(trends):
        x = 0.8 + i * 3.95
        add_card(slide, x, 2.0, 3.8, 4.5, COLOR_CARD)
        add_flat_icon(slide, x + 0.3, 2.2, 0.35, trend["icon"])
        add_text_box(slide, x + 0.2, 2.7, 3.4, 3.6, [
            {"text": trend["title"] + "\n\n", "font": "Arial", "size": 16, "bold": True, "color": COLOR_TITLE},
            {"text": trend["desc"], "font": "나눔고딕", "size": 13, "color": COLOR_TEXT}
        ])
        
    # Slide 29: 비즈니스 Action Plan - 기획/가격/마케팅/유통
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_BG)
    add_title(slide, "4대 핵심 비즈니스 액션 플랜")
    # 4분할 격자형 그리드 카드
    grid_plans = [
        {"title": "기획 부문 (AI 자동화 집중)", "desc": "단순한 챗봇 프롬프트를 넘어 '멀티 에이전트 오케스트레이션' 및 'n8n 노코드 자동화 연동' 등 고품질 실무 연동 기획에 투자", "x": 0.8, "y": 1.8},
        {"title": "가격 부문 (가치 포지셔닝)", "desc": "가격 무상관성(0.06)을 바탕으로 무리한 단가 인하 대신 고급 예제와 최신 버전을 보장하여 24,000~28,000원의 고단가 보전 채택", "x": 6.8, "y": 1.8},
        {"title": "마케팅 부문 (리뷰 개수 올인)", "desc": "리뷰 건수의 강한 상관성(0.48)을 고려하여 출간 초기 1개월 내 서평단 및 리워드로 임계점인 '누적 리뷰 50~100건' 조기 돌파", "x": 0.8, "y": 4.5},
        {"title": "유통 부문 (채널 차별화)", "desc": "양적 다각화는 한빛미디어의 유통망을, 질적 메가 히트를 노리는 차별화 타이틀은 이지스 혹은 골든래빗의 커뮤니티 채널 타겟팅", "x": 6.8, "y": 4.5}
    ]
    for plan in grid_plans:
        add_card(slide, plan["x"], plan["y"], 5.7, 2.3, COLOR_CARD)
        add_flat_icon(slide, plan["x"] + 0.3, plan["y"] + 0.2, 0.35, "diamond")
        add_text_box(slide, plan["x"] + 0.8, plan["y"] + 0.2, 4.7, 1.9, [
            {"text": plan["title"] + "\n", "font": "Arial", "size": 16, "bold": True, "color": COLOR_TITLE},
            {"text": plan["desc"], "font": "나눔고딕", "size": 12, "color": COLOR_TEXT}
        ])
        
    # Slide 30: 종합 결론 및 마무리
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, COLOR_LEAD_BG)
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    run1 = p1.add_run()
    run1.text = "경청해 주셔서 감사합니다"
    run1.font.name = 'Arial'
    run1.font.size = Pt(44)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_LEAD_TEXT
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(14)
    run2 = p2.add_run()
    run2.text = "본 분석 제언을 바탕으로 한 출판 성공률 극대화를 기원합니다."
    run2.font.name = '나눔고딕'
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0xD9, 0xE1, 0xF2)
    
    # ----------------------------------------------------
    # 각 슬라이드에 슬라이드 노트(발표 대본) 할당
    # ----------------------------------------------------
    print("슬라이드 노트(발표 대본) 주입 중...")
    for idx, slide in enumerate(prs.slides):
        if idx < len(scripts):
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = scripts[idx]
            
    # 프레젠테이션 저장
    print(f"PPTX 파일 저장 중: {output_filepath}")
    prs.save(output_filepath)
    print("PPTX 프레젠테이션 생성 완료!")

if __name__ == "__main__":
    dest_path = os.path.join("yes24", "docs", "eda_presentation.pptx")
    create_presentation(dest_path)
